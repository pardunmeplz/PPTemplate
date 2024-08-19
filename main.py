from pptx import Presentation
from pptx.util import Inches
import math
import os
import io
from pptx.enum.shapes import MSO_SHAPE, MSO_SHAPE_TYPE
from heic import convert_heic_to_jpeg_in_memory

# Function to extract slide layout properties
def get_slide_layout_properties(pptx_path, slide_number):
    prs = Presentation(pptx_path)
    
    # Get slide size
    slide_width = prs.slide_width
    slide_height = prs.slide_height
    
    # Get the specified slide (indexing starts at 0)
    slide = prs.slides[slide_number]
    
    positions = []
    
    # Loop through all shapes on the slide and find the ones that are images
    for shape in slide.shapes:
        if hasattr(shape, 'image'):
            # Get the position and size of the image
            left_inch = shape.left.inches
            top_inch = shape.top.inches
            width_inch = shape.width.inches
            height_inch = shape.height.inches
            
            # Append position and size to the list
            positions.append((left_inch, top_inch, width_inch, height_inch))
    
    # Return slide properties and positions
    return slide_width, slide_height, positions

def copy_text_frame(source_tf, dest_tf):

        dest_tf.margin_right = source_tf.margin_right
        dest_tf.margin_top = source_tf.margin_top
        dest_tf.margin_bottom = source_tf.margin_bottom
        dest_tf.margin_left = source_tf.margin_left
        dest_tf.vertical_anchor = source_tf.vertical_anchor
        dest_tf.word_wrap = source_tf.word_wrap
        dest_tf.auto_size = source_tf.auto_size

        for para_s, para_d in zip(source_tf.paragraphs, dest_tf.paragraphs):
            para_d.alignment = para_s.alignment
            para_d.level = para_s.level
            para_d.font.size = para_s.font.size
            para_d.font.name = para_s.font.name
            para_d.font.bold = para_s.font.bold
            para_d.font.italic = para_s.font.italic
            para_d.font.underline = para_s.font.underline
            
            
            for run_s in para_s.runs:
                run_d = para_d.add_run()
                run_d.text = run_s.text
                run_d.font.size = run_s.font.size
                run_d.font.name = run_s.font.name
                run_d.font.bold = run_s.font.bold
                run_d.font.italic = run_s.font.italic
                run_d.font.underline = run_s.font.underline

                if run_s.font.color and run_s.font.color.rgb:
                    run_d.font.color.rgb = run_s.font.color.rgb

def get_slide_copy(prs, slide_index, new_slide):
    
    source_slide = prs.slides[slide_index]

    # Copy all shapes (textboxes, images, lines, etc.) from the source slide to the new slide
    for shape in source_slide.shapes:
        if shape.has_text_frame:
            # Add a new textbox shape in the destination slide
            textbox = new_slide.shapes.add_textbox(shape.left, shape.top, shape.width, shape.height)
            new_text_frame = textbox.text_frame
            new_text_frame.clear()  # Clear any existing default text
            
            copy_text_frame(shape.text_frame, new_text_frame)

        elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            image_stream = io.BytesIO(shape.image.blob)
            new_slide.shapes.add_picture(image_stream, shape.left, shape.top, shape.width, shape.height)

        elif shape.shape_type == MSO_SHAPE_TYPE.LINE:
            new_shape = new_slide.shapes.add_shape(
                MSO_SHAPE_TYPE.LINE, shape.left, shape.top, shape.width, shape.height
            )
            new_shape.line.color.rgb = shape.line.color.rgb
            new_shape.line.width = shape.line.width

        elif shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
            new_shape = new_slide.shapes.add_shape(
                shape.auto_shape_type, shape.left, shape.top, shape.width, shape.height
            )
            new_shape.fill.fore_color.rgb = shape.fill.fore_color.rgb
            new_shape.line.color.rgb = shape.line.color.rgb

        else:
            print(f"Unsupported shape type: {shape.shape_type}. Skipping.")

# Function to open existing presentation or create a new one
def open_or_create_presentation(output_path, slide_width, slide_height):
    if os.path.exists(output_path):
        prs = Presentation(output_path)  # Open the existing presentation
        print(f"Appending to existing presentation: {output_path}")
    else:
        prs = Presentation()  # Create a new presentation
        prs.slide_width = slide_width
        prs.slide_height = slide_height
        print(f"Creating a new presentation: {output_path}")
    return prs

# Function to remove all placeholder shapes from the slide
def remove_placeholders(slide):
    for shape in slide.shapes:
        if shape.is_placeholder:
            sp = shape._element
            sp.getparent().remove(sp)

# Function to replace placeholders in the slide text
def replace_placeholders(slide, title_text=None, subtitle_text=None):
    for shape in slide.shapes:
        if shape.has_text_frame:
            text_frame = shape.text_frame
            for paragraph in text_frame.paragraphs:
                if title_text and "#title" in paragraph.text:
                    paragraph.text = paragraph.text.replace("#title", title_text)
                if subtitle_text and "#subtitle" in paragraph.text:
                    paragraph.text = paragraph.text.replace("#subtitle", subtitle_text)

# Function to add a title slide based on layout
def add_title_slide(prs, title_text):

    prs_layout = Presentation(layout_slide_path)    
    slide = create_blank_slide(prs)
    
    get_slide_copy(prs_layout, title_number, slide )
    
    # Replace placeholder text with actual title
    replace_placeholders(slide, title_text=title_text)

# Function to add a subtitle slide based on layout
def add_subtitle_slide(prs, title_text, subtitle_text):
    
    prs_layout = Presentation(layout_slide_path)
    slide = create_blank_slide(prs)

    get_slide_copy(prs_layout, subtitle_number, slide )
    
    # Replace placeholder text with actual title and subtitle
    replace_placeholders(slide, title_text=title_text, subtitle_text=subtitle_text)

# Function to create a truly blank slide
def create_blank_slide(prs):
    # Add a blank slide using the built-in blank layout
    slide_layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    
    # Remove all default placeholders
    remove_placeholders(slide)
    
    return slide

# Function to add images to a slide based on layout
def add_image_slide(prs, positions, image_files):
    slide = create_blank_slide(prs)  # Create a blank slide
    remove_placeholders(slide)  # Ensure no placeholders are on the slide

    # Place the images according to the layout template
    for img_path, pos in zip(image_files, positions):
        if '.HEIC' in img_path: img_path = convert_heic_to_jpeg_in_memory(img_path)
        left_inch, top_inch, width_inch, height_inch = pos
        slide.shapes.add_picture(img_path, left=left_inch * 914400, top=top_inch * 914400, 
                                 width=width_inch * 914400, height=height_inch * 914400)

# Recursive function to handle nested subfolders
def create_slides_for_folder(prs, layout_slide_path, slide_number, folder_path, parent_title=None):
    # Get image positions and layout properties from the layout slide
    _, _, positions = get_slide_layout_properties(layout_slide_path, slide_number)
    
    # Get the list of subfolders and files in the current folder
    entries = os.listdir(folder_path)
    subfolders = [os.path.join(folder_path, entry) for entry in entries if os.path.isdir(os.path.join(folder_path, entry))]
    image_files = [os.path.join(folder_path, entry) for entry in entries if entry.lower().endswith(('png', 'jpg', 'jpeg', 'heic'))]
    non_support_files = [os.path.join(folder_path, entry) for entry in entries 
                         if not entry.lower().endswith(('png', 'jpg', 'jpeg', 'heic')) and not os.path.isdir(os.path.join(folder_path, entry))]

    if len(non_support_files)> 0 :print("Following Files are not supported!!")
    for path in non_support_files:
        print(path)
    print()
    
    # Debug: Print number of image files found
    print(f"Found {len(image_files)} image(s) in {folder_path}")
    
    # Add image slides if there are images in this folder
    if image_files:
        images_per_slide = [image_files[x*len(positions):(x+1)*len(positions)] for x in range(math.ceil(len(image_files)/len(positions)))]
        for curr_image_files in images_per_slide:
            add_image_slide(prs, positions, curr_image_files)
    
    # Recursively handle subfolders
    for subfolder_path in subfolders:
        subfolder_name = os.path.basename(subfolder_path)
        
        # Add a title slide for the main subfolder
        if parent_title:
            add_subtitle_slide(prs, parent_title, subfolder_name)
        else:
            add_title_slide(prs, subfolder_name)
        
        # Recursively process the subfolder and its contents
        create_slides_for_folder(prs, layout_slide_path, slide_number, subfolder_path, parent_title=subfolder_name)

# Function to add images and subfolder slides
def create_presentation_with_nested_folders(layout_slide_path, slide_number, main_image_folder, output_path, title_number,subtitle_number):
    # Get image positions and layout properties from the layout slide
    slide_width, slide_height, positions = get_slide_layout_properties(layout_slide_path, slide_number)
    
    # Open existing presentation or create a new one
    prs = open_or_create_presentation(output_path, slide_width, slide_height)
    
    # Iterate over all first-level subfolders in the main image folder
    for folder_name in os.listdir(main_image_folder):
        folder_path = os.path.join(main_image_folder, folder_name)
        
        if os.path.isdir(folder_path):
            # Add a title slide for the main subfolder
            add_title_slide(prs, folder_name)
            
            # Recursively process the subfolder and its contents
            create_slides_for_folder(prs, layout_slide_path, slide_number, folder_path, parent_title=folder_name)
    
    # Save the updated or new presentation
    prs.save(output_path)
    print(f"Presentation saved: {output_path}")

# Paths and configuration
layout_slide_path = "input.pptx"  # Input PPTX with the layout
slide_number = 2  # Layout slide index in the template file
title_number = 0  # Title slide index in the template file
subtitle_number = 1  # subtitle slide index in the template file
main_image_folder = "images"  # Folder containing subfolders with images and nested subfolders
output_path = "output.pptx"  # Output file path for the generated/updated presentation

# Create or append to the presentation
create_presentation_with_nested_folders(layout_slide_path, slide_number, main_image_folder, output_path, title_number, subtitle_number)
